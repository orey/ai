import os, shutil, sys, csv
import subprocess
from pypdf import PdfReader
from pathlib import Path

from docx import Document
from pptx import Presentation
from openpyxl import load_workbook

import sys
sys.path.append('.')
from tools10 import CSV_trace, interrupt, create_text_file, generateName, Timer, is_zip, myprint, ensureFile
from manage_docx import get_document_text, treat_text_with_contents


SOURCE = "C:\\Users\\a876246\\Documents\\" # ends with \
TARGET = "C:\\c\\" # ends with \

ROOT = "oreyboulot-NHI" # Root folder will be added to source and target

LIBREOFFICE = "LibreOfficePortable"

#---------------------------------------------------------- Global traces
TRACE_HEADER = ["SOURCE_FILE_NAME","TARGET_FILE_NAME","STATUS"]

FILE_ALREADY_PROCESSED = "File already processed"

PDF_TEXT_EXTRACTED = "PDF text extracted"
PDF_OCR_DONE = "PDF OCR done"
PDF_NOT_READABLE = "PDF not readable"

DOCX_TEXT_EXTRACTED = "DOCX text extracted"
DOCX_NOT_READABLE = "DOCX not readable"

PPTX_TEXT_EXTRACTED = "PPTX text extracted"
PPTX_NOT_READABLE = "PPTX not readable"

XLSX_TEXT_EXTRACTED = "XLSX text extracted"
XLSX_NOT_READABLE = "XLSX not readable"

TEMP = generateName()
FILES_OK = CSV_trace(TEMP + "files_ok.csv")
FILES_OK.add(TRACE_HEADER)

# Some spreadsheets have worksheets with a number of rows that is superior
# to this limit. We decided to cap to ROWS_LIMIT
ROWS_LIMIT = 5000

# one column, no header
NOT_OK_FILENAME = "files_not_ok.csv"
FILES_NOT_OK = CSV_trace(NOT_OK_FILENAME, "a")


#---------------------------------------------------------- create_folder_if_not_exists
def create_folder_if_not_exists(path):
    """
    Create a folder if it doesn't exist; do nothing if it does.
    """
    try:
        thepath = Path(path)
        if thepath.exists() and thepath.is_dir():
            # Do nothing
            return True
        else:
            os.makedirs(path, exist_ok=True)
            myprint(f"| Folder | New folder created: {path}")
            return True
    except Exception as e:
        myprint(e)
        myprint("| Folder | Path may be too long. Skipping...")
        return False

    
#---------------------------------------------------------- copy_file
def copy_file(src, dst):
    """Copy src to dst only if dst does not exist."""
    if not os.path.exists(dst):
        shutil.copy2(src, dst)

        
#---------------------------------------------------------- remove_file
def remove_file(path):
    """Remove a file if it exists."""
    if os.path.exists(path):
        os.remove(path)

#---------------------------------------------------------- analyze_pdf
def get_text_from_standard_pdf(source_file, target_file, verbose=False):
    """
    Analyze a PDF to determine if it has extractable text or is image-only.
    Returns True if the PDF has extractable text, False otherwise.
    """
    try:
        reader = PdfReader(source_file)

        text = ""
        for page_num in range(len(reader.pages)):
            page = reader.pages[page_num]        
            text += page.extract_text()
    except Exception as e:
        myprint(f"| PDF  | Error | Error for file '{source_file}'")
        myprint(f"| PDF  | Exception type: {type(e).__name__}")
        myprint(f"| PDF  | Exception message: {e}")
        return False

    # If any non-whitespace text is found, consider it a standard PDF
    if len(text.strip())<2:
        myprint(f"| PDF  | No text found (image PDF?) | {source_file}", verbose)
        return False
    else:
        myprint(f"| PDF  | Standard PDF | {source_file} ", verbose)
    #creating output file
    return create_text_file(target_file + ".txt", text, prefix="| PDF  | ")


#---------------------------------------------------------- run_ocrmypdf
def run_ocrmypdf(input_pdf, output_pdf, language="eng", extra_args=None):
    """Run ocrmypdf on an input PDF and save the result to output_pdf."""
    cmd = ["ocrmypdf",
           "--force-ocr",
           "--sidecar",
           "-l", language,
           input_pdf,
           output_pdf]
    
    # Add any extra arguments if provided
    if extra_args:
        cmd.extend(extra_args)
    try:
        subprocess.run(cmd,
                       stdout=subprocess.DEVNULL,  # Suppress standard output
                       stderr=subprocess.DEVNULL,  # Suppress error output
                       check=True)
        myprint(f"| PDF  | {input_pdf} | OCR completed successfully: {output_pdf}")
        return True
    except subprocess.CalledProcessError as e:
        myprint(f"| PDF  | {input_pdf} | OCR failed with error code {e.returncode}")
        myprint(f"| PDF  | Skipping...")
        return False
        

#---------------------------------------------------------- process_pdf
def treat_pdf(source_file, target_file, verbose=False):
    """
    Process a PDF:
    - If it has extractable text, do nothing (standard PDF)
    - If it's image-only, run ocrmypdf to add OCR layer
    """
    
    # 1. Attempting to extract the text from PDF
    has_text = get_text_from_standard_pdf(source_file,target_file, verbose)

    if has_text:
        FILES_OK.add([source_file, target_file, PDF_TEXT_EXTRACTED])
        return True

    # 2. Attempting an OCR
    myprint(f"| PDF  | Attemting force scan | {source_file}")
    myprint(f"| PDF  | OCR started")
    if run_ocrmypdf(source_file, target_file, language="eng"):
        remove_file(target_file)
        FILES_OK.add([source_file, target_file, PDF_OCR_DONE])
        myprint("| PDF  | OCR ended OK")
        return True
    else:
        FILES_NOT_OK.add([source_file])
        myprint("| PDF  | PCR ended NOT OK")
        return False

    
#---------------------------------------------------------- treat_docx
def treat_docx(source_file, target_file, verbose=False):
    try:
        if not is_zip(source_file):
            myprint(f"| DOCX | Not a Zip | File {source_file}")
            FILES_NOT_OK.add([source_file])
            return False
        text = get_document_text(source_file)
        text2 = treat_text_with_contents(text)
        if create_text_file(target_file, text2, verbose=verbose, prefix="| DOCX | "):
            FILES_OK.add([source_file, target_file, DOCX_TEXT_EXTRACTED])
            return True
        else:
            FILES_NOT_OK.add([source_file])
            return False
    except Exception as e:
        myprint(e)
        FILES_NOT_OK.add([source_file])
        return False

    
#---------------------------------------------------------- treat_pptx
def treat_pptx(source_file, target_file, verbose=False):
    """
    Extract all text from a PowerPoint file.
    There is another more complex solution to extract text by
    extracting progressively all the items and metadata
    """
    try:
        if not is_zip(source_file):
            myprint(f"| PPTX | Not a Zip | File {source_file}")
            FILES_NOT_OK.add([source_file])
            return False
        prs = Presentation(source_file)
        text = ""
    
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text += shape.text + "\n"
        if create_text_file(target_file, text, verbose=verbose, prefix="| PPTX | "):
            FILES_OK.add([source_file, target_file, PPTX_TEXT_EXTRACTED])
            return True
        else:
            FILES_NOT_OK.add([source_file])
            return False
    except Exception as e:
        myprint(e)
        FILES_NOT_OK.add([source_file])
        return False


#---------------------------------------------------------- extract xlsx
def treat_xlsx(source_file, target_file, verbose=False):
    '''
    This method keeps all the empty lines.
    I don't know what it gives on merges cells.
    It can probably extract the pricing lists quite efficiently
    '''
    try:
        if not is_zip(source_file):
            myprint(f"| XLSX | Not a Zip | File {source_file}")
            FILES_NOT_OK.add([source_file])
            return False
        wb = load_workbook(source_file, read_only=True, data_only=True)
        sheets_status = {} # { sheet_name, True/False }

        for sheet_name in wb.sheetnames:
            # see if we have already processed the sheet
            target = target_file.replace(".xlsx","") + "_()_" + sheet_name + ".xlsx.txt"
            if ensureFile(target):
                myprint(f"| XLSX | File already processed: {target}")
                continue
            text = ""
            ws = wb[sheet_name]
            #text += f"Sheet named '{sheet_name}' content hereafter:\n\n"

            myprint(f"| XLSX | ()=> Sheet: '{sheet_name}' with rows/columns = {ws.max_row}/{ws.max_column}", verbose)
            if ws.max_row > ROWS_LIMIT:
                myprint(f"| XSLX | ()=> Sheet: '{sheet_name}' too large. Only {ROWS_LIMIT} rows will be processed", verbose)
            # main loop on rows
            countrows = 0
            for row in ws.iter_rows():
                countrows += 1
                if countrows > ROWS_LIMIT:
                    myprint(f"| XLSX | === Sheet: '{sheet_name}' too large. We have processed the {ROWS_LIMIT} rows only", verbose)
                    break
                if verbose:
                    print(".", end="", flush=True)
                row_data = []
                for cell in row:
                    # Skip empty cells if needed
                    value = cell.value if cell.value is not None else ""
                    row_data.append(str(value))

                # Process each row (myprint, save, etc.)
                #myprint(row_data)
                text += "| " + " | ".join(row_data) + ' |' + '\n'
            print("")
            #creating one file per spreadsheet
            if create_text_file(target, text, verbose=verbose, prefix="| XLSX | "):
                myprint(f"| XLSX | ()=> File for {sheet_name} created: {target}")
                sheets_status[sheet_name] = True
            else:
                sheets_status[sheet_name] = False
        wb.close()

        if all(sheets_status.values()):
            FILES_OK.add([source_file, target_file, XLSX_TEXT_EXTRACTED])
            return True
        else:
            FILES_NOT_OK.add([source_file])
            return False
    except Exception as e:
        myprint(e)
        return ""


#---------------------------------------------------------- main
def main():
    """
    Use os.walk to explore all files and folders in a directory.
    """
    
    # 1. Create the folders
    for root, dirs, files in os.walk(SOURCE + ROOT):
        #--- folder management
        for dir_name in dirs:
            thedir = os.path.join(root, dir_name)
            #myprint(f"[DIR] {thedir}")
            target_dir_name = os.path.join(root.replace(SOURCE,TARGET), dir_name)
            create_folder_if_not_exists(target_dir_name)

    # 2. Load a potential error file
    errorfiles = []
    with open(NOT_OK_FILENAME, "r", encoding="utf8") as f:
        reader = csv.reader(f, delimiter = ";")
        count = 0
        for line in reader:
            if count == 0:
                # we are in the header
                count += 1
                continue
            errorfiles.append(line[0])

    # 3. Manage the files
    count = {
        'txt': 0,
        'pdf': 0,
        'docx': 0,
        'pptx': 0,
        'xlsx': 0
    } # per {extension : extension}
    for root, dirs, files in os.walk(SOURCE + ROOT):
        for file_name in files:
            thefile = os.path.join(root, file_name)
            if thefile in errorfiles:
                # it is already in the error files, we can skip the treatment
                continue
            #myprint(f"[FILE] {thefile}")
            extension = file_name.split('.')[-1]
            target_file_name =  os.path.join(root.replace(SOURCE,TARGET),file_name)
            #--- TXT
            if extension.upper() == "TXT":
                count['txt'] += 1
                # the file name does not change
                if os.path.exists(target_file_name):
                    myprint(f"| TXT  | File already processed: {target_file_name}")
                    FILES_OK.add([thefile, target_file_name, FILE_ALREADY_PROCESSED])
                else:
                    copy_file(thefile, target_file_name)
                    myprint(f"| TXT  | The file {file_name} has been copied")
            #--- PDF
            elif extension.upper() == "PDF":
                count['pdf'] += 1
                # the real filename has a .txt extension
                if os.path.exists(target_file_name + ".txt"):
                    myprint(f"| PDF  | File already processed: {target_file_name}")
                    FILES_OK.add([thefile, target_file_name, FILE_ALREADY_PROCESSED])
                else:
                    # Note: there is no '.txt' in target because ocrmypdf
                    # needs a real name of PDF and adds the .txt itself.
                    # In case the standard extract works, the function
                    # adds it.
                    treat_pdf(thefile, target_file_name, True)
            #--- DOCX
            elif extension.upper() == "DOCX":
                count['docx'] += 1
                if os.path.exists(target_file_name + ".txt"):
                    myprint(f"| DOCX | File already processed: {target_file_name}")
                    FILES_OK.add([thefile, target_file_name, FILE_ALREADY_PROCESSED])
                else:
                    treat_docx(thefile, target_file_name + ".txt", verbose= True)
                    #interrupt(f"Look at {target_file_name}")
            #--- PPTX
            elif extension.upper() == "PPTX":
                count['pptx'] += 1
                if os.path.exists(target_file_name + ".txt"):
                    myprint(f"| PPTX | File already processed: {target_file_name}")
                    FILES_OK.add([thefile, target_file_name, FILE_ALREADY_PROCESSED])
                else:
                    treat_pptx(thefile, target_file_name + ".txt", verbose= True)
                    #interrupt(f"Look at {target_file_name}")
            #--- XLSX
            elif extension.upper() == "XLSX":
                count['xlsx'] += 1
                if os.path.exists(target_file_name + ".txt"):
                    myprint(f"| XLSX | File already processed: {target_file_name}")
                    FILES_OK.add([thefile, target_file_name, FILE_ALREADY_PROCESSED])
                else:
                    # no extension is given for target file because we have
                    # one file per sheet.
                    treat_xlsx(thefile, target_file_name, verbose= True)
                    #interrupt(f"Look at {target_file_name}")
            #--- everything else
            else:
                if extension in count:
                    count[extension] += 1
                else:
                    count[extension] = 1
    myprint(f"==> {len(count)} files visited\nSplit per type:")
    for elem in count:
        myprint(f"| {elem} | {count[elem]} files")


#======================================= entry point
if __name__ == "__main__":
    t = Timer("main")
    main()
    t.stop()


